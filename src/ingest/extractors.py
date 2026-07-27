"""Extractores de contenido por formato.

Cada extractor devuelve una lista de dicts: {"text": str, "metadata": {...}}.
Hoy solo se usa `extract_markdown` (los 16 documentos fuente son .md), pero se
deja el dispatcher completo porque `SUPPORTED_EXTENSIONS` en config.py ya
contempla otros formatos para cuando se sumen documentos reales (PDF, DOCX, etc.).
"""

import csv
import json
import re
from pathlib import Path
from typing import Any, Dict, List

from loguru import logger


def extract_pdf(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de PDF nativo usando pymupdf4llm (salida Markdown, página por página)."""
    import pymupdf4llm

    try:
        pages = pymupdf4llm.to_markdown(str(file_path), page_chunks=True)
        return [
            {
                "text": page.get("text", ""),
                "metadata": {
                    "page": page.get("metadata", {}).get("page", 0) + 1,
                    "format": "pdf",
                    "extraction_method": "pymupdf4llm",
                },
            }
            for page in pages
        ]
    except Exception as e:
        logger.warning(f"pymupdf4llm falló para {file_path}, intentando OCR: {e}")
        return extract_pdf_ocr(file_path)


def extract_pdf_ocr(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de PDF escaneado usando OCR (Tesseract). Fallback de extract_pdf."""
    import io

    import fitz
    import pytesseract
    from PIL import Image

    doc = fitz.open(str(file_path))
    results = []

    for page_num, page in enumerate(doc):
        pix = page.get_pixmap(dpi=300)
        img = Image.open(io.BytesIO(pix.tobytes("png")))
        text = pytesseract.image_to_string(img, lang="spa+eng")

        if text.strip():
            results.append(
                {
                    "text": text,
                    "metadata": {
                        "page": page_num + 1,
                        "format": "pdf_ocr",
                        "extraction_method": "tesseract",
                    },
                }
            )

    doc.close()
    return results


def extract_docx(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de Word preservando estructura de headings como sección."""
    from docx import Document

    doc = Document(str(file_path))
    results = []
    current_section = "Inicio"

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        if para.style.name.startswith("Heading"):
            current_section = text

        results.append(
            {
                "text": text,
                "metadata": {
                    "section": current_section,
                    "style": para.style.name,
                    "format": "docx",
                    "extraction_method": "python-docx",
                },
            }
        )

    return results


def extract_xlsx(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae contenido de Excel como texto estructurado, una entrada por hoja."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    results = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        headers = [str(cell) if cell else f"Col_{i}" for i, cell in enumerate(rows[0])]
        text_lines = [f"## Hoja: {sheet_name}", f"Columnas: {', '.join(headers)}", ""]

        for row in rows[1:]:
            row_text = " | ".join(
                f"{headers[i]}: {cell}" for i, cell in enumerate(row) if cell is not None
            )
            if row_text.strip():
                text_lines.append(row_text)

        results.append(
            {
                "text": "\n".join(text_lines),
                "metadata": {"sheet": sheet_name, "format": "xlsx", "extraction_method": "openpyxl"},
            }
        )

    wb.close()
    return results


def extract_pptx(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de PowerPoint incluyendo notas del orador."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    results = []

    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"[Notas del orador]: {notes}")

        if texts:
            results.append(
                {
                    "text": "\n".join(texts),
                    "metadata": {
                        "slide": slide_num,
                        "has_notes": bool(notes),
                        "format": "pptx",
                        "extraction_method": "python-pptx",
                    },
                }
            )

    return results


HEADING_PATTERN = re.compile(r"^(#{1,6})\s+(.+)$")


def extract_markdown(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de Markdown segmentado por encabezado.

    Cada segmento queda etiquetado con su heading en `metadata["section"]`,
    requerido por yachay-rag-pipeline para citar archivo + sección de origen.
    """
    text = file_path.read_text(encoding="utf-8")
    lines = text.split("\n")

    segments: List[Dict[str, Any]] = []
    current_section = "Inicio"
    current_lines: List[str] = []

    def flush():
        body = "\n".join(current_lines).strip()
        if body:
            segments.append(
                {
                    "text": body,
                    "metadata": {
                        "section": current_section,
                        "format": "markdown",
                        "extraction_method": "native",
                    },
                }
            )

    for line in lines:
        match = HEADING_PATTERN.match(line.strip())
        if match:
            flush()
            current_section = match.group(2).strip()
            current_lines = [line]
        else:
            current_lines.append(line)
    flush()

    if not segments:
        segments = [
            {
                "text": text,
                "metadata": {"section": "Documento completo", "format": "markdown", "extraction_method": "native"},
            }
        ]

    return segments


def extract_csv_file(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae CSV como texto estructurado (una línea por fila)."""
    text_lines = []
    with open(file_path, "r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        headers = reader.fieldnames or []
        text_lines.append(f"Columnas: {', '.join(headers)}")
        for row in reader:
            row_text = " | ".join(f"{k}: {v}" for k, v in row.items() if v)
            text_lines.append(row_text)

    return [{"text": "\n".join(text_lines), "metadata": {"format": "csv", "extraction_method": "csv.DictReader"}}]


def extract_json_file(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae JSON como texto legible (pretty-printed)."""
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    text = json.dumps(data, indent=2, ensure_ascii=False)
    return [{"text": text, "metadata": {"format": "json", "extraction_method": "json.load"}}]


def extract_html(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto de HTML eliminando tags."""
    from bs4 import BeautifulSoup

    html = file_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")
    text = soup.get_text(separator="\n", strip=True)
    return [{"text": text, "metadata": {"format": "html", "extraction_method": "beautifulsoup4"}}]


def extract_text(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae texto plano."""
    text = file_path.read_text(encoding="utf-8")
    return [{"text": text, "metadata": {"format": "txt", "extraction_method": "native"}}]


EXTRACTOR_MAP = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".md": extract_markdown,
    ".csv": extract_csv_file,
    ".json": extract_json_file,
    ".html": extract_html,
    ".txt": extract_text,
}


def extract_file(file_path: Path) -> List[Dict[str, Any]]:
    """Extrae contenido de cualquier formato soportado, delegando al extractor correspondiente."""
    ext = file_path.suffix.lower()
    extractor = EXTRACTOR_MAP.get(ext)
    if not extractor:
        logger.warning(f"Formato no soportado: {ext} ({file_path.name})")
        return []

    logger.info(f"Extrayendo: {file_path.name} (formato: {ext})")
    results = extractor(file_path)

    for r in results:
        r["metadata"]["source_file"] = file_path.name
        r["metadata"]["source_path"] = str(file_path)

    return results
